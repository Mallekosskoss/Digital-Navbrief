from pptx import Presentation
import pandas as pd

# Define the path to your PowerPoint file
pptx_file = '../data_inn/Navbrief MPN5 Vestland dag 1.pptx'

# Define the title of the slide from which you want to extract tables
target_slide_title = 'Navigasjonssystem- Instillinger og hensyn'

# Load the PowerPoint presentation
presentation = Presentation(pptx_file)

# Find the slide with the specified title
target_slide = None
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text == target_slide_title:
            target_slide = slide
            break

if target_slide is None:
    print(f"Slide with title '{target_slide_title}' not found.")
    exit()

# Extract tables from the target slide
table_data = []
for shape in target_slide.shapes:
    if shape.has_table:
        table = shape.table
        table_data.append(table)

# Create a Pandas DataFrame to preserve the table structure
df_list = []
for table in table_data:
    df = pd.DataFrame(columns=range(len(table.columns)))
    for row in table.rows:
        data = [cell.text for cell in row.cells]
        df.loc[len(df)] = data
    df_list.append(df)

# Print each table in a structured format
for idx, df in enumerate(df_list):
    print(f"Table {idx + 1}:\n")
    print(df.to_string(index=False))
    print("\n")